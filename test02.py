import os
import pandas as pd
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm,Pt,Cm
from pptx.dml.color import RGBColor
from openpyxl import Workbook,load_workbook
from shutil import copyfile,copytree
import pypinyin


k='圣诞老人来了'

kk=pypinyin.pinyin(k)

kkk=[c[0] for c in kk]

print(' '.join(kkk))

# ppt='I:\\乐高\\图纸\\template2.pptx'

# df=pd.read_excel('E:\\WXWork\\1688852895928129\\WeDrive\\大智小超科学实验室\\2-乐高课程\\课程信息表.xlsx')

# t1=df[df['课程编号']=='L046']['知识点'].values.tolist()[0]
# t1=t1[2:]

# t2=df[df['课程编号']=='L046']['课后问题'].values.tolist()[0]
# t2=t2[2:]

# ptn=r"\d."
# _t1=re.sub(ptn,'',t1)
# _t2=re.sub(ptn,'',t2)

# print(_t1)



# prs=Presentation(ppt)  
# blank_slide_layout=prs.slide_layouts[1]



# prs.slides[2].placeholders[1].text=_t1
# slide=prs.slides.add_slide(blank_slide_layout)
# prs.slides[len(prs.slides)-1].placeholders[0].text='课后问题'
# prs.slides[len(prs.slides)-1].placeholders[1].text=_t2

# slide=prs.slides.add_slide(blank_slide_layout)
# picTitle=slide.shapes.add_picture('e:\\temp\\end_ppt_pic.png',Cm(0),Cm(1.4),height=Cm(17.69)) #加入封面图
# slide.shapes._spTree.insert(2, picTitle._element)

# prs.slides[4].shapes.title.text='kkkkkkoooooo'
# # print(prs.slides[2].placeholders[1].text)

# prs.save('e:\\temp\\kkkk.pptx')