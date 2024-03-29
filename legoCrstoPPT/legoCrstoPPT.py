import os
import sys
# sys.path.append('i:/py/dzxc/module')
sys.path.append(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.realpath(__file__)))),'modules'))
from readConfig import readConfig
from PIL import Image,ImageDraw,ImageFont
import platform
import numpy as np
import re
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm,Pt,Cm
from pptx.dml.color import RGBColor
from openpyxl import Workbook,load_workbook
from shutil import copyfile,copytree
import pypinyin

class picToPPT:
    def __init__(self,crsName):
        self.crsName=crsName
        config=readConfig(os.path.join(os.path.dirname(os.path.realpath(__file__)),'legoCrstoPPT.config'))

        if platform.system().lower()=='linux':
            self.picDir='/home/jack/data/乐高/图纸'
            self.template='/home/jack/data/乐高/图纸/template.pptx'
            self.wtmark='/home/jack/data/乐高/图纸/大智小超水印.png'
            self.crs_info=''
        else:
            self.picDir=config['图纸文件夹']
            self.template=config['PPT模板']
            self.wtmark=config['PPT水印']
            self.legoCodeList=config['wedo零件编号']
            self.crs_info=config['课程信息表']
            self.blk_pic_dir=config['零件图片文件夹']
            self.box_pos_pic=os.path.join(config['盒子位置图文件夹'],'盒子位置编号图.jpg')
        self.picSrc=os.path.join(self.picDir,crsName)
        # self.block_names=self.blockNames()

    def makeDirs(self,dirName):
        if not os.path.exists(dirName):
            print('新项目，正在创建文件夹……',end='')
            os.makedirs(dirName)
            os.makedirs(os.path.join(dirName,'零件总图'))
            os.makedirs(os.path.join(dirName,'video'))
            os.makedirs(os.path.join(dirName,'animation'))
            print('创建完成,请在文件夹中放入步骤图片。')
            os.startfile(dirName)
            sys.exit(0)
        else:
            print('已有项目……',end='')
            files=os.listdir(dirName)
            if len(files)<3:
                print('请先做好步骤图片')
                sys.exit (0)

    def test_stepXls(self,lxfml_mode):
        xlsName=os.path.join(self.picSrc,self.crsName+'-ppt步骤零件名称.xlsm')
        if not os.path.exists(xlsName):
            blNames=self.blockNames(mode=lxfml_mode)
            # blNames=self.block_names
            copyfile(os.path.join(self.picDir,'ppt_step_Template.xlsm'),xlsName)
            wb=load_workbook(xlsName,keep_vba=True)
            tb=wb.active
            tb['A1']='序号'
            tb['B1']='零件名称'
            tb['D1']='零件名称（识别）'
            tb['E1']='零件位置'
            row=2
            for file in os.listdir(self.picSrc):
                ptn='\d{3}_\dx.png'                
                if re.match(ptn,file):
                    tb['A'+str(row)]=row-1
                    row+=1
            
            for k,blk in enumerate(blNames):
                tb['D'+str(k+2)]=blk[0]
                tb['E'+str(k+2)]=blk[1]
                
            wb.save(xlsName)
            print('新建了步骤文件，请在文件中先输入零件名称。')
            sys.exit(0)
        else:
            df=pd.read_excel(xlsName,usecols=[0,1])
            blockNum=df['零件名称'].count()
            if blockNum==0:
                print('未输入零件名称，请先写入。')
                sys.exit(0)

    def renameFiles(self):
        print('正在查看是否需要修改文件名……',end='')
        for file in os.listdir(self.picSrc):
            if file[-3:].lower()=='png':
                if len(file)<10:
                    oldname=os.path.join(self.picSrc,file)
                    newname=os.path.join(self.picSrc,file.zfill(10))
                    os.rename(oldname,newname)
        print('完成')

    def copytoCrsDir(self,crsPPTDir='I:\\乐高\\乐高WeDo\\课程'):
        desDir=os.path.join(crsPPTDir,self.crsName)
        if not os.path.exists(desDir):
            os.makedirs(desDir)
        
        oldName_video=os.path.join(self.picSrc,'video')
        oldName_ppt=os.path.join(self.picSrc,self.crsName+'_00.pptx')
        newName_video=os.path.join(desDir,'video')
        newName_ppt=os.path.join(desDir,self.crsName+'.pptx')

        copyfile(oldName_ppt,newName_ppt)
        if not os.path.exists(newName_video):
            copytree(oldName_video,newName_video)
        os.startfile(desDir)
        os.startfile(newName_ppt)

    def blockNames(self,mode='new',add_block='no',add_list=''):
        fn=os.path.join(self.picDir,self.crsName,self.crsName[4:]+'.lxfml')
        bl=pd.read_excel(self.legoCodeList)
        bl['lxfml编号'].astype('object')
        with open (fn,'r',encoding='utf-8') as f:
            rl=f.readlines()

        if mode=='old':
            ptn1='(?<=\<Part refID\=\"\d{1}\" designID\=\")(.*)(?=\;A\" materials\=)'
            ptn2='(?<=\<Part refID\=\"\d{2}\" designID\=\")(.*)(?=\;A\" materials\=)'
            ptn3='(?<=\<Part refID\=\"\d{3}\" designID\=\")(.*)(?=\" partType\=)'
        elif mode=='new':
            #studio更新的至2.2.2(1)版本后，lxfml文档结构改变
            ptn1='(?<=\<Part refID\=\"\d{1}\" designID\=\")(.*)(?=\" partType\=)'
            ptn2='(?<=\<Part refID\=\"\d{2}\" designID\=\")(.*)(?=\" partType\=)'
            ptn3='(?<=\<Part refID\=\"\d{3}\" designID\=\")(.*)(?=\" partType\=)'
        else:
            print('lxfml文件的内容不能识别')
            exit()

        a=[]
        for line in rl:
            p=re.findall(ptn1,line)
            if p:
                a.append(p[0])

            p=re.findall(ptn2,line)
            if p:
                a.append(p[0])

            p=re.findall(ptn3,line)
            if p:
                a.append(p[0])
                
        out=[]
        for  aa in a:
            try:
                aa=int(aa)
            except:
                pass
            try:
                blkName=bl[bl['lxfml编号']==aa]['lxfml显示名称'].values.tolist()[0]
                blk_pos=bl[bl['lxfml编号']==aa]['位置'].values.tolist()[0]
                blk_pic_addr=bl[bl['lxfml编号']==aa]['图片文件名'].values.tolist()[0]
                out.append([blkName,blk_pos,blk_pic_addr])
            except:
                if aa!=19079:
                    out.append([aa,'-'])
        
        
        #如有表中未列出的新零件，临时增加
        if add_block=='yes':
            # print(add_list)
            for add_blk_id in add_list:
                for rpt in range(add_blk_id[1]):
                    # print(add_blk_id,rpt)
                    to_add=[bl[bl['lxfml编号']==add_blk_id[0]]['lxfml显示名称'].values.tolist()[0],
                            bl[bl['lxfml编号']==add_blk_id[0]]['位置'].values.tolist()[0],
                            bl[bl['lxfml编号']==add_blk_id[0]]['图片文件名'].values.tolist()[0]]
                    out.append(to_add)

        # print('out',out)
        # self.block_names=out
        # out['图片文件名'].astype(str).apply(lambda x : os.path.join(self.blk_pic_dir,x))
        # out['图片文件名']=os.path.join(self.blk_pic_dir,out['图片文件名'].astype(str))
        # print(out[2])
        return out

    def block_pic_list(self,save='yes',lxfml_mode='new',add_block='no',add_list=''):
        block_names_file=self.blockNames(mode=lxfml_mode,add_block=add_block,add_list=add_list)
        df=pd.DataFrame(block_names_file)
        df.columns=['零件名称','位置','图片地址']
        #如无图片地址则丢掉，此行有可能导致计算数量错误
        df.dropna(axis=0,how='any',subset=['图片地址'],inplace=True)
        df['图片地址']=df['图片地址'].apply(lambda x:os.path.join(self.blk_pic_dir,x))
        df_gp=pd.DataFrame(df.groupby('零件名称').count())
        df_gp.reset_index(inplace=True)
        adr=df.drop_duplicates(['零件名称'])
        dict_adr=dict(zip(adr['零件名称'],adr['图片地址']))
        dict_pos=dict(zip(adr['零件名称'],adr['位置']))
        df_gp['图片地址']=df_gp['零件名称'].apply(lambda x: dict_adr[x])
        df_gp.columns=['零件名称','数量','图片地址']
        df_gp['位置']=df_gp['零件名称'].apply(lambda x: dict_pos[x])

        df_gp.sort_values(by=['位置'],ascending=False,inplace=True)
        df_gp.reset_index(inplace=True)
        
        df_box=df.groupby('位置').count()
        df_box.reset_index(inplace=True)
        df_box.columns=['位置','数量','图片地址']
        # print(df_box)


        h_pic=200
        gap=10
        bg_h=(h_pic+gap)*df_gp.shape[0]+300
        bg_w=800
        bg=Image.new('RGBA',[bg_w,bg_h],'#FFFFFF')
        draw=ImageDraw.Draw(bg)
        posxy=[100,250]     

        txt_title=self.crsName+'  零件分类清单'
        # print(txt_title)
        # print((bg_w-len(txt_title)*40)//2)
        # draw.text([100,80],txt_title,'#7EB554',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',40))
        draw.text([(bg_w-len(txt_title)*40)//2,80],txt_title,'#7EB554',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',40))
        draw.text([(bg_w-len(txt_title)*40)//2,160],'总数量： '+str(df_gp['数量'].sum()),'#8FAE8E',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',32))

        for index,row in df_gp.iterrows():

            pic=self.pic_resize(pic_adr=row['图片地址'],h=h_pic)
            bg.paste(pic[0],posxy,mask=pic[1])      
            draw.text([posxy[0]+270,posxy[1]+70],row['零件名称'],'#AE9D8E',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',30))      
            draw.text([posxy[0]+270,posxy[1]+70+50],'×'+str(row['数量']),'#FF0000',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',36))      
            if index>0 :
                #不是最后一行
                if index<(df_gp.shape[0]-1):
                    if df_gp.iloc[index]['位置']!=df_gp.iloc[index-1]['位置']:
                        draw.line([(0,posxy[1]-gap//2),(bg_w,posxy[1]-gap//2)],'#000000')  
                        draw.text([bg_w-200,posxy[1]-gap*8],str(df_gp.iloc[index-1]['位置']),'#6E6D6D',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',60))   
                        draw.text([bg_w-190,posxy[1]-gap*20],'×'+str(df_box[df_box['位置']==df_gp.iloc[index-1]['位置']]['数量'].tolist()[0]),'#F490EF',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',86))  
                #最后一行
                else:
                    if df_gp.iloc[index]['位置']==df_gp.iloc[index-1]['位置']:
                        draw.text([bg_w-200,bg_h-200],str(df_gp.iloc[index-1]['位置']),'#6E6D6D',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',60))
                        draw.text([bg_w-190,posxy[1]-gap*8],'×'+str(df_box[df_box['位置']==df_gp.iloc[index-1]['位置']]['数量'].tolist()[0]),'#F490EF',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',86))  
                    else:
                        draw.line([(0,posxy[1]-gap//2),(bg_w,posxy[1]-gap//2)],'#000000')  
                        draw.text([bg_w-200,posxy[1]-gap*8],str(df_gp.iloc[index-1]['位置']),'#6E6D6D',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',60))  
                        draw.text([bg_w-200,bg_h-200],str(df_gp.iloc[index-1]['位置']),'#6E6D6D',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',60))
                        draw.text([bg_w-190,posxy[1]-gap*20],'×'+str(df_box[df_box['位置']==df_gp.iloc[index-1]['位置']]['数量'].tolist()[0]),'#F490EF',font=ImageFont.truetype('c:\\windows\\fonts\\simhei.ttf',86))  
            
            posxy=[posxy[0],posxy[1]+h_pic+gap]

        if save=='yes':
            savename=os.path.join(self.picSrc,self.crsName+'_零件分类清单.jpg')
            bg=bg.convert('RGB')
            bg.save(savename)
            print('零件分类清单已经保存至 {}'.format(savename))
        else:
            print('仅显示，未保存。')
            bg.show()
    

    def pic_resize(self,pic_adr,h=480):
        img=Image.open(pic_adr)
        img=img.resize((h*img.size[0]//img.size[1],h))
        # print(pic_adr,img.split())
        return img,img.split()[3]

    def inner_box_pos(self,save='yes',lxfml_mode='new'):
        block_names_file=self.blockNames(mode=lxfml_mode)
        box_pic_jpg=self.box_pos_pic
        df=pd.DataFrame(block_names_file)

        df.columns=['零件名称','位置','图片地址']
        df.dropna(axis=0,how='any',subset=['图片地址'])
        gp=df.groupby(['位置']).count()
        pos=list(gp.index)
        values=[]
        for v in gp.values:
            values.append(v[0])
        df_pos=pd.DataFrame({'位置':pos,'数量':values})
        font_size_pos_1=200
        font_size_pos_2=150
        pos_data={
            '0-0':[2086,2432,font_size_pos_1],
            '1-1':[3032,1737,font_size_pos_1],
            '1-2':[3032,1023,font_size_pos_1],
            '1-3':[3032,276,font_size_pos_1],
            '2-1':[2432,1494,font_size_pos_2],
            '2-2':[2432,615,font_size_pos_2],
            '2-3':[2432,186,font_size_pos_2],
            '3-1':[2048,1494,font_size_pos_2],
            '3-2':[2048,615,font_size_pos_2],
            '3-3':[2048,186,font_size_pos_2],
            '4-1':[1340,1446,font_size_pos_1],
            '4-2':[1340,372,font_size_pos_1],
            '5-1':[434,1683,font_size_pos_1],
            '5-2':[434,486,font_size_pos_1]  
        }

        bg=Image.open(box_pic_jpg)
        draw=ImageDraw.Draw(bg)
        pos_dat=df_pos.loc[df_pos['位置']!='-']

        for pos_num in pos_dat['位置'].tolist():
            # if pos_num!='0-0':
                # print(df_pos[df_pos['位置']==pos_num]['数量'].values)
            draw.text((pos_data[pos_num][0],pos_data[pos_num][1]),str(df_pos[df_pos['位置']==pos_num]['数量'].values[0]),fill='red',font=ImageFont.truetype('j:\\fonts\\FZYunDongCuHei.ttf',pos_data[pos_num][2]))
            draw.text((100,2462),'总数：'+str(pos_dat['数量'].sum()),fill='#6074b7',font=ImageFont.truetype('j:\\fonts\\FZYunDongCuHei.ttf',110))

        if save=='yes':
            savename=os.path.join(self.picSrc,self.crsName+'_零件位置图.jpg')
            bg.save(savename)
            print('零件位置图已经保存至 {}'.format(savename))
        else:
            print('仅显示，未保存。')
            bg.show()

        

    def ExpPPT(self,copyToCrsDir='no',lxfml_mode='new',crsPPTDir='I:\\乐高\\乐高WeDo\\课程'):
        print('\n正在处理：')   
        self.makeDirs(self.picSrc)
        self.renameFiles()    
        self.test_stepXls(lxfml_mode=lxfml_mode)    
        def picList():
            ptn_block_list=re.compile(r'\d*_1x_tagged.png')
            ptn_block_list2=re.compile(r'\d*_1x.png')
            for blk_old in os.listdir(self.picSrc):
                if blk_old[-3:].lower()=='png' and len(blk_old)<10:
                    os.rename(os.path.join(self.picSrc,blk_old),os.path.join(self.picSrc,blk_old.zfill(10)))
                
            for blk_old2 in os.listdir(self.picSrc):
                if blk_old2[-3:].lower()=='png' and len(blk_old2)<10:
                    os.rename(os.path.join(self.picSrc,'零件总图',blk_old2),os.path.join(self.picSrc,'零件总图',blk_old2.zfill(10)))
            
            pic_steps=[]
            for blk in os.listdir(os.path.join(self.picSrc,'零件总图')):
                # if ptn_block_list.match(blk) or ptn_block_list2.match(blk):
                if ptn_block_list2.match(blk): #找出零件总图的数量
                    pic_steps.append(os.path.join(self.picSrc,'零件总图',blk))
            pic_steps.sort()
            summary_num=len(pic_steps)
        
            ptn=r'\d.*png'
            picList=[]
            for fn in os.listdir(self.picSrc):
                if re.match(ptn,fn):
                    picList.append(os.path.join(self.picSrc,fn))
            picList.sort()
            #             print(picList)

            pic_steps.extend(picList)
    
            #             print(pic_steps)
    
            return [pic_steps,summary_num]

        def picToPPT(picList):
            #读取课程信息表中的知识点和课后问题
            df=pd.read_excel(self.crs_info)
            t1=df[df['课程编号']==self.crsName[0:4]]['知识点'].values.tolist()[0]
            t1=t1[2:]
            t2=df[df['课程编号']==self.crsName[0:4]]['课后问题'].values.tolist()[0]
            t2=t2[2:]

            #用空格替换原来的1. 2. 3. 格式
            ptn=r"\d."
            txt_knlg=re.sub(ptn,'',t1)
            txt_question=re.sub(ptn,'',t2)

            step_blkList=pd.read_excel(os.path.join(self.picSrc,self.crsName+'-ppt步骤零件名称.xlsm'),usecols=[0,1]).replace(np.nan,'')['零件名称'].tolist()
            #             print(step_blkList)
            
            
            prs=Presentation(self.template)            
            left=Cm(0)
            top=Cm(1.4)
            height=Cm(17.69)
            
            #             left_wtmk=Cm(5)
            #             top_wtmk=Cm(5)
            left_wtmk=Cm(0)
            top_wtmk=Cm(1.4)

            lowerPosList_1=['2x16单位黑色板','12单位绿色梁','12单位绿色孔砖','1×16梁']
            lowerPosList_2=['主机','电机']

            blank_slide_layout=prs.slide_layouts[1]

            slide=prs.slides.add_slide(blank_slide_layout)
            picTitle=slide.shapes.add_picture(os.path.join(self.picSrc,self.crsName[4:]+'.jpg'),left,top,height=height) #加入封面图
            slide.shapes._spTree.insert(2, picTitle._element)#将图片放在底层
            slide.shapes.title.text=self.crsName[4:] #课程名称 
            textbox = slide.shapes.add_textbox(Cm(3),Cm(6),Cm(5),Cm(3))
            tf = textbox.text_frame
            para = tf.add_paragraph()    # 新增段落
            _pinyin=[c[0] for c in pypinyin.pinyin(self.crsName[4:])] #写入标题拼音
            title_pinyin=' '.join(_pinyin)

            para.text = title_pinyin  # 向段落写入文字

            #写入零件数及步数
            slide=prs.slides.add_slide(blank_slide_layout)
            t_blk_num=df[df['课程编号']==self.crsName[0:4]]['零件数'].values.tolist()[0]
            t_step_num=df[df['课程编号']==self.crsName[0:4]]['步数'].values.tolist()[0]
            t_blk_step='完成搭建需要{}步\n一共需要{}个零件'.format(int(t_step_num),int(t_blk_num))
            slide.shapes.placeholders[0].text='步数和零件数'
            slide.shapes.placeholders[1].text_frame.text=t_blk_step 

            for pgh in slide.shapes.placeholders[1].text_frame.paragraphs:
                for run in pgh.runs:
                    run.font.size=Pt(40)


            for i,img in enumerate(picList[0]):                
                slide=prs.slides.add_slide(blank_slide_layout)
                if img[-3:].lower()=='png':
                    pic=slide.shapes.add_picture(img,left,top,height=height)
                
               
                if i>=picList[1]: #跳过零件总图的数量
                #                     print(i,i-picList[1])
                    try:                        
                        txt=step_blkList[i-picList[1]] #根据文件内容决定文本框的位置下移程度，以免遮挡零件图片
                        if txt in lowerPosList_1:
                            y_textbox=Cm(6)
                        elif txt in lowerPosList_2:
                            y_textbox=Cm(7)
                        else:
                            y_textbox=Cm(5)
                        textbox=slide.shapes.add_textbox(Cm(2),y_textbox,Cm(5),Cm(2.5))
                        p = textbox.text_frame.add_paragraph()
                        p.line_spacing=1.2 #行间距


                        p.text=txt
                        p.font.size=Pt(24)
                        p.font.color.rgb = RGBColor(22, 56, 153)
                    except:
                        pass
                
                pic_wtmark=slide.shapes.add_picture(self.wtmark,left_wtmk,top_wtmk) #加水印
                

            prs.slides[2].placeholders[1].text=txt_knlg #写入知识点
            slide=prs.slides.add_slide(blank_slide_layout)

            prs.slides[len(prs.slides)-1].placeholders[0].text='课堂分享'  #加入课堂分享
            prs.slides[len(prs.slides)-1].placeholders[1].text=txt_question
            for pgh in prs.slides[len(prs.slides)-1].placeholders[1].text_frame.paragraphs:
                for run in pgh.runs:    
                    run.font.size=Pt(32)

            slide=prs.slides.add_slide(blank_slide_layout)
            picTitle=slide.shapes.add_picture(os.path.join(self.picDir,'end_ppt_pic.png'),Cm(0),Cm(1.4),height=Cm(17.69)) #加入收拾零件的图片

            newFn=os.path.join(self.picSrc,self.crsName+'_00.pptx')
            prs.save(newFn)

            if copyToCrsDir=='yes':
                self.copytoCrsDir(crsPPTDir=crsPPTDir)

            print('ppt已导出完成，文件名：{}'.format(newFn))                              
            
        pic_list=picList()
        picToPPT(pic_list)          
        
if __name__=='__main__':
    mypics=picToPPT('L149可爱的招财猫')
    # mypics.inner_box_pos(save='yes',lxfml_mode='new')
    # print(mypics.blockNames())
    # k=mypics.blockNames()   
    # mypics=picToPPT('/home/jack/data/乐高/图纸/031回力赛车')
    # mypics.ExpPPT(copyToCrsDir='no',lxfml_mode='new',crsPPTDir='I:\\乐高\\乐高WeDo\\课程')
    mypics.block_pic_list(save='yes',lxfml_mode='new',add_block='no',add_list=[[85546,2]])
    # mypics.makeDirs()
    # mypics.copytoCrsDir()
